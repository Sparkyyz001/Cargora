# Генератор презентации для SmartScape Hackathon 2026.
# Выход: presentation/Cargora_SmartScape.pptx (9 слайдов, обязательная
# структура: проблема / решение / технология / демонстрация / перспективы).

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).parent
OUT = ROOT.parent / "presentation"
OUT.mkdir(exist_ok=True)

BG = RGBColor(0x0B, 0x12, 0x20)        # тёмный фон
FG = RGBColor(0xF1, 0xF5, 0xF9)        # основной текст
MUT = RGBColor(0x94, 0xA3, 0xB8)       # приглушённый
ACC = RGBColor(0x0E, 0xA5, 0xE9)       # акцент (sky-500, как в продукте)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def text(s, x, y, w, h, runs, size=18, color=FG, bold=False, align=PP_ALIGN.LEFT, spacing=1.0):
    """runs: строка или список (текст, dict-переопределения)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for t, opt in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", align)
        p.line_spacing = spacing
        p.space_after = Pt(opt.get("space_after", 6))
        r = p.add_run()
        r.text = t
        f = r.font
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.color.rgb = opt.get("color", color)
        f.name = "Segoe UI"
    return box


def accent_bar(s, y=Inches(1.35), x=Inches(0.7), w=Inches(0.9)):
    bar = s.shapes.add_shape(1, x, y, w, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACC
    bar.line.fill.background()


def header(s, kicker, title):
    text(s, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.4), kicker, size=13, color=ACC, bold=True)
    text(s, Inches(0.7), Inches(0.65), Inches(11.9), Inches(0.7), title, size=30, bold=True)
    accent_bar(s)


def bullets(s, items, x=Inches(0.7), y=Inches(1.7), w=Inches(11.9), size=17):
    runs = []
    for head, rest in items:
        runs.append((f"•  {head}", {"size": size, "bold": True, "space_after": 2}))
        if rest:
            runs.append((f"    {rest}", {"size": size - 2, "color": MUT, "space_after": 12}))
    text(s, x, y, w, Inches(5.4), runs)


# ── 1. Титул ─────────────────────────────────────────────────────────────────
s = slide()
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
     "SmartScape Hackathon 2026  ·  Трек 1: Smart Mobility & Infrastructure",
     size=15, color=ACC, bold=True)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.3), "Cargora", size=66, bold=True)
text(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(1.0),
     "AI-платформа Каспийского транзитного коридора:\nпрогноз загруженности пунктов пропуска на 48 часов вперёд",
     size=22, color=MUT)
text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.9), [
    ("Живое демо: cargora.vercel.app", {"size": 18, "color": FG, "bold": True}),
    ("Команда: «________»", {"size": 16, "color": MUT}),
])

# ── 2. Проблема ──────────────────────────────────────────────────────────────
s = slide()
header(s, "ПРОБЛЕМА", "Транзит через Каспий стоит в очередях")
bullets(s, [
    ("Средний коридор перегружен", "после 2022 года грузопоток через Казахстан вырос кратно; порт Актау и КПП работают на пределе"),
    ("Фуры стоят на КПП «Болашак» сутками", "перевозчик узнаёт об очереди, только когда уже приехал на границу"),
    ("Шторм на Каспии останавливает порт", "ветер >13 м/с прекращает погрузку — очередь судов копится непредсказуемо"),
    ("Нет предиктивной информации", "существующие сервисы показывают очередь «сейчас», но никто не прогнозирует «завтра»"),
    ("Цена вопроса", "простой фуры — $250–400/сутки; для города — потерянная пропускная способность коридора"),
])

# ── 3. Решение ───────────────────────────────────────────────────────────────
s = slide()
header(s, "РЕШЕНИЕ", "Cargora — предиктивная диспетчеризация транзита")
bullets(s, [
    ("ML-прогноз загруженности на 48 часов", "собственная модель прогнозирует загрузку порта Актау, КПП «Болашак», «Тажен» и ж/д узла Бейнеу"),
    ("Рекомендация оптимального окна прохождения", "система подсказывает 4-часовое окно с минимальной очередью — груз не стоит, КПП разгружается"),
    ("ИИ-подбор рейса с учётом прогноза", "заявка → оптимальный паром или автоперевозчик + обоснование, когда привезти груз в порт"),
    ("Единая платформа", "диспетчерская в реальном времени (Supabase Realtime), трекинг на карте 2GIS, аналитика транзита для акимата"),
    ("Охват", "Казахстан · Туркменистан · Азербайджан · Иран · Россия — 5 стран одним договором"),
])

# ── 4. Технология ────────────────────────────────────────────────────────────
s = slide()
header(s, "ТЕХНОЛОГИЯ", "Собственная ML-модель, а не готовое API")
text(s, Inches(0.7), Inches(1.65), Inches(6.4), Inches(4.8), [
    ("Пайплайн", {"size": 18, "bold": True, "space_after": 8}),
    ("1. Симулятор транзитного потока — датасет 71 040 наблюдений (2 года, 4 пункта): суточный/недельный/годовой профили, погода, праздники РК, инциденты", {"size": 14, "color": MUT, "space_after": 8}),
    ("2. GradientBoostingRegressor (scikit-learn): 220 деревьев, валидация временным сплитом без утечки", {"size": 14, "color": MUT, "space_after": 8}),
    ("3. Экспорт деревьев в JSON → инференс на TypeScript внутри Next.js — без Python-сервера в проде", {"size": 14, "color": MUT, "space_after": 8}),
    ("4. Вход модели — реальный прогноз погоды Open-Meteo; LLM (Groq) генерирует обоснование поверх нашей модели", {"size": 14, "color": MUT, "space_after": 14}),
    ("Метрики (тест: март–июнь 2026)", {"size": 18, "bold": True, "space_after": 8}),
    ("MAE 4.04%  ·  RMSE 5.81  ·  R² 0.896", {"size": 16, "color": ACC, "bold": True, "space_after": 4}),
    ("Бейзлайн «час недели»: MAE 6.85% — наша модель точнее на 41%", {"size": 14, "color": MUT}),
])
s.shapes.add_picture(str(ROOT / "figures" / "feature_importance.png"), Inches(7.3), Inches(1.8), width=Inches(5.5))

# ── 5. Качество модели ───────────────────────────────────────────────────────
s = slide()
header(s, "ТЕХНОЛОГИЯ", "Прогноз против факта — тестовое окно 14 дней")
s.shapes.add_picture(str(ROOT / "figures" / "forecast_vs_actual.png"), Inches(1.9), Inches(1.55), height=Inches(5.5))
text(s, Inches(0.7), Inches(7.0), Inches(11.9), Inches(0.4),
     "Выброс 8–9 июня на КПП «Тажен» — случайный инцидент (латентный фактор): модель честно не переобучена под шум",
     size=12, color=MUT)

# ── 6. Демонстрация ──────────────────────────────────────────────────────────
s = slide()
header(s, "ДЕМОНСТРАЦИЯ", "Живой продукт: cargora.vercel.app")
img = ROOT.parent / "public" / "dashboard-preview.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(4.4), Inches(1.7), width=Inches(8.2))
text(s, Inches(0.7), Inches(1.75), Inches(3.6), Inches(5.0), [
    ("Что показываем вживую:", {"size": 16, "bold": True, "space_after": 10}),
    ("• Заявка → ИИ подбирает рейс и окно прибытия в порт", {"size": 13, "color": MUT, "space_after": 8}),
    ("• GET /api/forecast — прогноз 4 пунктов на 48 ч", {"size": 13, "color": MUT, "space_after": 8}),
    ("• Диспетчерская: заявка прилетает в реальном времени", {"size": 13, "color": MUT, "space_after": 8}),
    ("• Трекинг груза на карте 2GIS", {"size": 13, "color": MUT, "space_after": 8}),
    ("• Аналитика транзита региона для акимата", {"size": 13, "color": MUT, "space_after": 8}),
])

# ── 7. Практическая применимость ─────────────────────────────────────────────
s = slide()
header(s, "ПРИМЕНИМОСТЬ", "Кому это нужно уже сегодня")
bullets(s, [
    ("Перевозчики и грузоотправители", "планируют прибытие в окно с минимальной очередью — экономия $250–400 на каждые сутки простоя фуры"),
    ("Порт Актау и КПП", "сглаживание пиков нагрузки: равномерный поток вместо штурма в часы пик"),
    ("Акимат Мангистауской области", "региональная аналитика транзита: объёмы, динамика, структура грузов — для решений по инфраструктуре"),
    ("Готовность к продакшену", "развёрнуто на Vercel, БД с RLS-политиками, ролевая модель, подписочная монетизация"),
    ("Масштабирование", "модель переобучается на реальных данных КГД/порта за один запуск train.py; добавление нового КПП — одна строка конфигурации"),
])

# ── 8. Перспективы ───────────────────────────────────────────────────────────
s = slide()
header(s, "ПЕРСПЕКТИВЫ", "Дорожная карта")
bullets(s, [
    ("Реальные данные", "интеграция с КГД РК и АО «Порт Актау»: обучение на фактических логах пропуска вместо симулятора"),
    ("Больше пунктов", "весь Средний коридор: Курык, Достык, Хоргос, Сарыагаш — модель уже мультипунктовая"),
    ("Уведомления", "push водителю: «выезжай в 02:00 — пройдёшь границу без очереди»"),
    ("Прогноз цен", "вторая модель: динамика ставок фрахта по коридору"),
    ("B2G-кабинет", "панель для акимата и Минтранса: прогноз нагрузки коридора на неделю вперёд"),
])

# ── 9. Финал ─────────────────────────────────────────────────────────────────
s = slide()
text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2), "Cargora", size=54, bold=True)
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.8),
     "Груз не должен стоять в очереди, если очередь можно предсказать.",
     size=22, color=ACC, bold=True)
text(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.4), [
    ("Демо: cargora.vercel.app", {"size": 18, "bold": True}),
    ("GitHub: github.com/Sparkyyz001/Cargora", {"size": 16, "color": MUT}),
    ("Трек 1: Smart Mobility & Infrastructure  ·  SmartScape Hackathon 2026", {"size": 14, "color": MUT}),
])

path = OUT / "Cargora_SmartScape.pptx"
prs.save(path)
print(f"OK: {path}")
